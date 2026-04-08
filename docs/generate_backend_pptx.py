"""견적 원가 분석 시스템 - 백엔드 설계서 PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

RED = RGBColor(230, 0, 18)
DARK = RGBColor(10, 22, 40)
BLUE = RGBColor(0, 56, 117)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 100)
LGRAY = RGBColor(245, 247, 250)
BLACK = RGBColor(25, 31, 40)
GREEN = RGBColor(46, 125, 50)
ORANGE = RGBColor(230, 81, 0)
PURPLE = RGBColor(123, 31, 162)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def slide(title, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    b.fill.solid(); b.fill.fore_color.rgb = RED; b.line.fill.background()
    b2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3))
    b2.fill.solid(); b2.fill.fore_color.rgb = DARK; b2.line.fill.background()
    ft = b2.text_frame; ft.text = "현대모비스 견적 원가 분석 시스템  |  백엔드 설계서  |  Confidential"
    ft.paragraphs[0].font.size = Pt(8); ft.paragraphs[0].font.color.rgb = RGBColor(150,150,150); ft.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tx = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
    p = tx.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = DARK
    if sub:
        p2 = tx.text_frame.add_paragraph(); p2.text = sub; p2.font.size = Pt(14); p2.font.color.rgb = GRAY
    return s

def text(s, l, t, w, h, txt, sz=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return tf

def mtext(s, l, t, w, h, lines, sz=11, color=BLACK):
    tx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]; p.font.bold = line[1]; p.font.size = Pt(line[2] if len(line)>2 else sz)
            p.font.color.rgb = line[3] if len(line)>3 else color
        else:
            p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color
        p.space_after = Pt(3)
    return tf

def box(s, l, t, w, h, fill=LGRAY, border=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    return sh

def table(s, l, t, w, data, cw=None):
    r, c = len(data), len(data[0])
    ts = s.shapes.add_table(r, c, Inches(l), Inches(t), Inches(w), Inches(0.32*r))
    tb = ts.table
    if cw:
        for i, ww in enumerate(cw): tb.columns[i].width = Inches(ww)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = BLACK
                if ri == 0: p.font.bold = True; p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if ri == 0 else (WHITE if ri%2==1 else LGRAY)
    return ts

# ═══════════════════════════════════════
# 표지
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(0.06))
b.fill.solid(); b.fill.fore_color.rgb = RED; b.line.fill.background()
text(s, 1.5, 1.0, 10, 1, "견적 원가 분석 시스템", 44, True, WHITE, PP_ALIGN.CENTER)
text(s, 1.5, 2.0, 10, 0.6, "백엔드 설계서", 28, False, RGBColor(200,200,200), PP_ALIGN.CENTER)
text(s, 1.5, 3.5, 10, 0.5, "Spring Boot 3.2  |  PostgreSQL 16  |  Redis 7  |  Nginx 1.25", 16, False, RGBColor(180,180,180), PP_ALIGN.CENTER)
text(s, 1.5, 4.2, 10, 0.5, "현대모비스  |  원가관리팀  |  v2.0  |  2026.04.08", 14, False, RGBColor(140,140,140), PP_ALIGN.CENTER)

# 목차
s = slide("목차", "Table of Contents")
mtext(s, 0.8, 1.2, 11, 5, [
    "1.  시스템 아키텍처",
    "2.  기술 스택",
    "3.  패키지 구조",
    "4.  API 엔드포인트 — 인증 / 사용자",
    "5.  API 엔드포인트 — 견적서 파싱 / 검증",
    "6.  API 엔드포인트 — 분석 / 비교",
    "7.  API 엔드포인트 — 모델관리 / 변경요청",
    "8.  API 엔드포인트 — 인사이트 / 알림 / 대시보드",
    "9.  데이터베이스 ERD",
    "10. 주요 테이블 명세",
    "11. 인증/보안 설계",
    "12. 프론트엔드 ↔ 백엔드 매핑",
    "13. 배포 구성 (Docker Compose + Nginx)",
], 16, DARK)

# 1. 아키텍처
s = slide("1. 시스템 아키텍처", "System Architecture")
comps = [
    ("Browser\n(React)", 4.5, 0.3, 4, 0.8, RGBColor(0,100,255)),
    ("Nginx :80\n정적파일 + 리버스 프록시", 4.5, 1.6, 4, 0.8, RGBColor(46,125,50)),
    ("Spring Boot :8080\nREST API (WAS)", 4.5, 3.0, 4, 0.8, RED),
    ("PostgreSQL\n:5432", 0.8, 4.5, 2.5, 0.8, RGBColor(51,103,145)),
    ("Redis\n:6379", 4.2, 4.5, 2, 0.8, RGBColor(220,53,69)),
    ("MinIO\n:9000", 7, 4.5, 2, 0.8, ORANGE),
    ("Claude API\n(AI)", 10, 4.5, 2.5, 0.8, PURPLE),
]
for label, x, y, w, h, color in comps:
    box(s, x, y+1, w, h, WHITE, color)
    text(s, x+0.1, y+1.05, w-0.2, h-0.1, label, 11, True, color, PP_ALIGN.CENTER)
for y1, y2 in [(1.8, 2.4), (2.8, 3.4), (3.8, 4.8)]:
    text(s, 6.3, y1+0.7, 0.4, 0.4, "↓", 18, True, GRAY, PP_ALIGN.CENTER)

# 2. 기술 스택
s = slide("2. 기술 스택", "Technology Stack")
table(s, 0.6, 1.3, 12, [
    ["구분", "기술", "버전", "용도"],
    ["WAS", "Spring Boot", "3.2.4", "REST API 서버"],
    ["언어", "Java", "17 (LTS)", "백엔드 개발"],
    ["ORM", "Spring Data JPA", "-", "데이터 접근"],
    ["DB", "PostgreSQL", "16", "주 데이터 저장소 (JSONB 지원)"],
    ["캐시", "Redis", "7", "JWT 블랙리스트, 세션 캐시"],
    ["파일저장소", "MinIO", "latest", "Excel/PDF 원본 파일 (S3 호환)"],
    ["웹서버", "Nginx", "1.25", "리버스 프록시, 정적 파일, gzip"],
    ["인증", "JWT (jjwt)", "0.12.5", "Access/Refresh Token"],
    ["문서", "SpringDoc OpenAPI", "2.4.0", "Swagger UI 자동 생성"],
    ["Excel", "Apache POI", "5.2.5", "Excel 파싱"],
    ["AI", "Claude API", "claude-opus-4-6", "인사이트 채팅"],
    ["컨테이너", "Docker Compose", "3.9", "로컬/운영 배포"],
], [1.5, 2.5, 1.5, 6.5])

# 3. 패키지 구조
s = slide("3. 패키지 구조", "com.costanalysis.*")
box(s, 0.6, 1.3, 5.8, 5.5, WHITE, RGBColor(200,200,200))
text(s, 0.8, 1.4, 5.5, 0.4, "domain/ (비즈니스 도메인)", 14, True, BLUE)
mtext(s, 0.8, 1.9, 5.5, 4.7, [
    ("auth/", True, 11, DARK),   "  controller + dto + service",
    ("user/", True, 11, DARK),   "  controller + dto + entity + repository + service",
    ("quotation/", True, 11, DARK), "  controller(3) + dto(5) + entity(3) + repo(3) + service(5)",
    ("verification/", True, 11, DARK), "  controller + dto(2) + entity + repo + service",
    ("analysis/", True, 11, DARK), "  controller + dto + service",
    ("comparison/", True, 11, DARK), "  controller + dto(2) + entity + repo + service",
    ("model/", True, 11, DARK),  "  controller(2) + dto(5) + entity(2) + repo(2) + service(2)",
    ("insight/", True, 11, DARK), "  controller + dto(4) + entity(2) + repo(2) + service(2)",
    ("notification/", True, 11, DARK), "  controller + dto + entity(2) + repo(2) + service",
    ("dashboard/", True, 11, DARK), "  controller + dto + service",
], 10)

box(s, 6.8, 1.3, 5.8, 5.5, WHITE, RGBColor(200,200,200))
text(s, 7, 1.4, 5.5, 0.4, "global/ (공통 인프라)", 14, True, BLUE)
mtext(s, 7, 1.9, 5.5, 4.7, [
    ("config/", True, 11, DARK),
    "  SecurityConfig.java — Spring Security",
    "  CorsConfig.java — CORS 설정",
    "  OpenApiConfig.java — Swagger",
    "  WebClientConfig.java — Claude API",
    "  MinioConfig.java — 파일 스토리지",
    "",
    ("security/", True, 11, DARK),
    "  JwtTokenProvider.java — 토큰 생성/검증",
    "  JwtAuthenticationFilter.java — 요청 인증",
    "",
    ("exception/", True, 11, DARK),
    "  ErrorCode.java — 에러 코드 enum",
    "  BusinessException.java — 비즈니스 예외",
    "  GlobalExceptionHandler.java — @RestControllerAdvice",
    "",
    ("response/", True, 11, DARK),
    "  ApiResponse.java — 공통 응답 { success, data, error }",
    "",
    ("storage/", True, 11, DARK),
    "  FileStorageService.java — MinIO 업/다운로드",
], 10)

# 4. API — 인증/사용자
s = slide("4. API 엔드포인트 — 인증 / 사용자", "Auth & User APIs")
table(s, 0.6, 1.3, 12, [
    ["Method", "Endpoint", "설명", "인증", "권한"],
    ["POST", "/api/v1/auth/login", "로그인 (JWT 발급)", "X", "-"],
    ["POST", "/api/v1/auth/refresh", "토큰 갱신", "Refresh", "-"],
    ["POST", "/api/v1/auth/logout", "로그아웃 (블랙리스트)", "O", "USER"],
    ["GET", "/api/v1/users/me", "내 프로필 조회", "O", "USER"],
    ["PUT", "/api/v1/users/me", "프로필 수정", "O", "USER"],
    ["PUT", "/api/v1/users/me/password", "비밀번호 변경", "O", "USER"],
    ["GET", "/api/v1/users", "전체 사용자 목록", "O", "ADMIN"],
    ["POST", "/api/v1/users", "사용자 생성", "O", "ADMIN"],
], [0.8, 3, 3, 1.5, 3.7])

# 5. API — 견적서/검증
s = slide("5. API 엔드포인트 — 견적서 파싱 / 검증", "Quotation & Verification APIs")
table(s, 0.6, 1.3, 12, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/api/v1/quotations", "견적서 목록 (필터/검색/정렬/페이징)"],
    ["GET", "/api/v1/quotations/{id}", "견적서 상세"],
    ["POST", "/api/v1/quotations/upload", "파일 업로드 (multipart, 최대 50MB)"],
    ["POST", "/api/v1/quotations/{id}/parse", "파싱 실행 (비동기)"],
    ["DELETE", "/api/v1/quotations/{id}", "견적서 삭제"],
    ["GET", "/api/v1/quotations/{id}/parsed-items", "파싱 항목 목록"],
    ["PUT", "/api/v1/quotations/{id}/status", "상태 변경"],
    ["GET", "/api/v1/quotations/notes?fileId={id}", "파일별 노트 목록"],
    ["POST", "/api/v1/quotations/notes", "노트 생성"],
], [0.8, 4.5, 6.7])

table(s, 0.6, 4.8, 12, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/api/v1/verifications/{quotationId}", "검증 결과 조회"],
    ["POST", "/api/v1/verifications/{quotationId}/auto", "자동 검증 실행"],
    ["PUT", "/api/v1/verifications/{id}/decide", "수동 검증 (승인/수정)"],
    ["PUT", "/api/v1/verifications/{quotationId}/complete", "검증 완료 처리"],
], [0.8, 4.5, 6.7])

# 6. API — 분석/비교
s = slide("6. API 엔드포인트 — 분석 / 비교", "Analysis & Comparison APIs")
table(s, 0.6, 1.3, 12, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/api/v1/analysis/{quotationId}", "분석 결과 조회 (3그룹 + 이상치)"],
    ["POST", "/api/v1/analysis/{quotationId}/run", "분석 실행"],
    ["GET", "/api/v1/analysis/{quotationId}/golden-set", "골든셋 조회"],
    ["GET", "/api/v1/analysis/{quotationId}/export/excel", "Excel 다운로드 (.xls)"],
], [0.8, 4.5, 6.7])

table(s, 0.6, 3.5, 12, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/api/v1/comparisons", "비교 세션 목록"],
    ["POST", "/api/v1/comparisons", "비교 세션 생성 (아이템 + 견적서 2~4개)"],
    ["GET", "/api/v1/comparisons/{id}", "비교 결과 조회 (나란히 비교)"],
    ["DELETE", "/api/v1/comparisons/{id}", "비교 세션 삭제"],
], [0.8, 4.5, 6.7])

# 7. API — 모델관리/변경요청
s = slide("7. API 엔드포인트 — 모델관리 / 변경 요청", "Cost Model & Change Request APIs")
table(s, 0.6, 1.3, 5.8, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/formulas", "수식 목록"],
    ["POST", "/formulas", "수식 추가 (ADMIN)"],
    ["PUT", "/formulas/{id}", "수식 수정 (ADMIN)"],
    ["DEL", "/formulas/{id}", "수식 삭제 (ADMIN)"],
], [0.6, 1.8, 3.4])
text(s, 0.6, 1.1, 5.8, 0.3, "Base: /api/v1/models", 10, False, GRAY)

table(s, 6.8, 1.3, 5.8, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/change-requests", "전체 목록"],
    ["GET", ".../formula/{id}", "수식별 목록"],
    ["POST", "/change-requests", "변경 요청 생성"],
    ["PUT", ".../{id}/approve", "승인 (부서 선택, ADMIN)"],
    ["PUT", ".../{id}/reject", "반려 (ADMIN)"],
    ["DEL", ".../{id}", "취소 (본인 PENDING)"],
], [0.6, 2, 3.2])
text(s, 6.8, 1.1, 5.8, 0.3, "Base: /api/v1/models", 10, False, GRAY)

box(s, 0.6, 4.2, 12, 2.8, WHITE, RGBColor(200,200,200))
text(s, 0.8, 4.3, 11.5, 0.4, "변경 요청 워크플로우", 14, True, BLUE)
mtext(s, 0.8, 4.8, 11.5, 2.0, [
    "1. 일반 사용자 → POST /change-requests (수정 내용 + 사유) → status: PENDING",
    "2. 관리자 → GET /change-requests (목록 조회, status 필터링)",
    "3-A. 승인 → PUT /{id}/approve { comment, departments: ['견적1팀','견적2팀'] } → 수식에 자동 반영",
    "3-B. 반려 → PUT /{id}/reject { comment: '반려 사유' }",
    "4. 취소 → DELETE /{id} (본인 PENDING 상태만 가능)",
], 11)

# 8. API — 인사이트/알림/대시보드
s = slide("8. API 엔드포인트 — 인사이트 / 알림 / 대시보드", "Insight & Notification & Dashboard APIs")
table(s, 0.6, 1.3, 5.8, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/sessions", "채팅 세션 목록"],
    ["POST", "/sessions", "새 세션 생성"],
    ["GET", "/sessions/{id}/messages", "메시지 조회"],
    ["POST", "/sessions/{id}/chat", "메시지 전송 (AI 응답)"],
    ["DEL", "/sessions/{id}", "세션 삭제"],
], [0.6, 2.5, 2.7])
text(s, 0.6, 1.1, 5.8, 0.3, "인사이트: /api/v1/insights", 10, False, GRAY)

table(s, 6.8, 1.3, 5.8, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/notifications", "알림 목록"],
    ["PUT", "/read-all", "모두 읽음"],
    ["GET", "/activities", "작업 이력 (타입 필터)"],
], [0.6, 2.2, 3])
text(s, 6.8, 1.1, 5.8, 0.3, "알림: /api/v1/notifications", 10, False, GRAY)

table(s, 0.6, 4.5, 12, [
    ["Method", "Endpoint", "설명"],
    ["GET", "/api/v1/dashboard/stats", "요약 통계 (총 견적서, 검증완료율, 이상치, 평균원가)"],
    ["GET", "/api/v1/dashboard/work-items", "내 작업 현황 (6개 상태별 카운트)"],
    ["GET", "/api/v1/dashboard/action-alerts", "업무 알림 (긴급/주의/정보)"],
], [0.8, 3.5, 7.7])

# 9. ERD
s = slide("9. 데이터베이스 ERD", "Entity Relationship Diagram — PostgreSQL 16")
tables_info = [
    ("users", 0.5, 1.3, BLUE),
    ("quotations", 3, 1.3, GREEN),
    ("parsed_items", 5.8, 1.3, GREEN),
    ("verification_results", 8.5, 1.3, ORANGE),
    ("parsing_notes", 11, 1.3, ORANGE),
    ("cost_formulas", 0.5, 3.5, PURPLE),
    ("change_requests", 3, 3.5, PURPLE),
    ("insight_sessions", 5.8, 3.5, RGBColor(0,100,255)),
    ("insight_messages", 8.5, 3.5, RGBColor(0,100,255)),
    ("notifications", 0.5, 5.2, RGBColor(220,53,69)),
    ("activity_logs", 3, 5.2, RGBColor(220,53,69)),
    ("comparison_sessions", 5.8, 5.2, RGBColor(156,39,176)),
]
for name, x, y, color in tables_info:
    box(s, x, y, 2.2, 0.7, WHITE, color)
    text(s, x+0.1, y+0.1, 2, 0.5, name, 10, True, color, PP_ALIGN.CENTER)

# 10. 테이블 명세
s = slide("10. 주요 테이블 명세", "Table Specifications")
table(s, 0.6, 1.3, 12, [
    ["테이블", "설명", "주요 컬럼", "예상 규모"],
    ["users", "사용자", "employee_id, name, dept, role(ADMIN/USER)", "~100"],
    ["quotations", "견적서 메타", "file_name, status(6단계), progress, parsed_items, anomaly_count", "~10,000"],
    ["parsed_items", "파싱된 원가 항목", "category, name, amount, confidence, status, cell_ref", "~500,000"],
    ["verification_results", "검증 결과", "original_value, parsed_value, corrected_value, status", "~500,000"],
    ["cost_formulas", "원가 수식", "name, expression, badge, departments(JSONB)", "~50"],
    ["change_requests", "변경 요청", "formula_id, status, original/modified(JSONB), approved_depts", "~500"],
    ["parsing_notes", "파싱 노트", "file_id, user_id, type, content", "~5,000"],
    ["insight_sessions", "AI 채팅 세션", "user_id, title", "~1,000"],
    ["insight_messages", "채팅 메시지", "session_id, role(USER/AI), content", "~10,000"],
    ["notifications", "알림", "user_id, title, type, read", "~50,000"],
    ["activity_logs", "작업 이력", "user_id, action, target_type, target_id", "~100,000"],
], [1.8, 1.5, 5.5, 3.2])

# 11. 인증/보안
s = slide("11. 인증/보안 설계", "JWT Authentication & RBAC")
box(s, 0.6, 1.3, 5.8, 5.5, WHITE, RGBColor(200,200,200))
text(s, 0.8, 1.4, 5.5, 0.4, "JWT 인증 흐름", 14, True, BLUE)
mtext(s, 0.8, 1.9, 5.5, 4.7, [
    ("1. 로그인", True, 12, DARK),
    "  POST /api/v1/auth/login { employeeId, password }",
    "  → Access Token (1시간) + Refresh Token (7일)",
    "",
    ("2. API 요청", True, 12, DARK),
    "  Authorization: Bearer {accessToken}",
    "  → JwtAuthenticationFilter 검증",
    "  → SecurityContext 사용자 설정",
    "",
    ("3. 토큰 갱신", True, 12, DARK),
    "  POST /api/v1/auth/refresh",
    "  → Refresh Token 검증 → 새 Access Token",
    "",
    ("4. 로그아웃", True, 12, DARK),
    "  POST /api/v1/auth/logout",
    "  → Redis 블랙리스트 등록",
    "",
    ("5. 계정 잠금", True, 12, RED),
    "  5회 연속 실패 → Redis 카운트 → 자동 잠금",
], 11)

box(s, 6.8, 1.3, 5.8, 5.5, WHITE, RGBColor(200,200,200))
text(s, 7, 1.4, 5.5, 0.4, "역할 기반 접근 제어 (RBAC)", 14, True, BLUE)
table(s, 7, 2.0, 5.5, [
    ["기능", "ADMIN", "USER"],
    ["파일 업로드/파싱", "O", "O"],
    ["데이터 검증", "O", "O"],
    ["원가 분석/비교", "O", "O"],
    ["Excel 다운로드", "O", "O"],
    ["수식 직접 편집/삭제", "O", "X"],
    ["변경 요청 승인/반려", "O", "X"],
    ["부서 적용 범위 설정", "O", "X"],
    ["수정/추가 요청 제출", "-", "O"],
    ["본인 요청 취소", "-", "O"],
], [2.5, 1.5, 1.5])

# 12. 프론트 ↔ 백엔드 매핑
s = slide("12. 프론트엔드 ↔ 백엔드 매핑", "Frontend-Backend Integration")
table(s, 0.6, 1.3, 12, [
    ["프론트 페이지", "경로", "백엔드 Controller", "API Base", "비고"],
    ["로그인", "/login", "AuthController", "/api/v1/auth", "JWT 발급"],
    ["대시보드", "/dashboard", "DashboardController", "/api/v1/dashboard", "집계 통계"],
    ["견적서 분석", "/parsing_card", "QuotationController", "/api/v1/quotations", "업로드+파싱"],
    ["데이터 검증", "/verification", "VerificationController", "/api/v1/verifications", "원본 대조"],
    ["원가 분석", "/analysis", "AnalysisController", "/api/v1/analysis", "4탭 뷰+골든셋"],
    ["견적서 비교", "/comparison", "ComparisonController", "/api/v1/comparisons", "3단계 비교"],
    ["모델관리", "/models", "CostModel+ChangeRequest", "/api/v1/models", "수식+변경요청"],
    ["인사이트", "/insight", "InsightController", "/api/v1/insights", "AI 채팅"],
    ["이력/알림", "/history", "NotificationController", "/api/v1/notifications", "이력+알림"],
    ["설정", "/settings", "UserController", "/api/v1/users/me", "프로필/비밀번호"],
], [1.5, 1.3, 2.5, 2.5, 4.2])

mtext(s, 0.6, 5.2, 12, 1.5, [
    ("localStorage → DB 전환", True, 13, BLUE),
    "  cost-analysis-formulas → cost_formulas | cost-analysis-change-requests → change_requests",
    "  cost-analysis-current-user → users (JWT) | parsing-notes → parsing_notes",
    "  guide-dismissed-{path} → 프론트 유지 (개인 UI 설정, DB 불필요)",
], 11)

# 13. 배포 구성
s = slide("13. 배포 구성", "Docker Compose + Nginx")
table(s, 0.6, 1.3, 12, [
    ["서비스", "이미지", "포트", "볼륨", "용도"],
    ["nginx", "nginx:1.25-alpine", "80", "nginx.conf + React build", "웹서버 + 리버스 프록시"],
    ["app", "Dockerfile (Java 17)", "8080", "-", "Spring Boot WAS"],
    ["postgres", "postgres:16-alpine", "5432", "postgres_data + schema.sql", "메인 DB"],
    ["redis", "redis:7-alpine", "6379", "redis_data", "캐시 + JWT 블랙리스트"],
    ["minio", "minio/minio:latest", "9000/9001", "minio_data", "파일 스토리지"],
], [1.2, 2.5, 1.2, 3.5, 3.6])

box(s, 0.6, 4.2, 5.8, 2.8, WHITE, RGBColor(200,200,200))
text(s, 0.8, 4.3, 5.5, 0.4, "실행 방법", 14, True, GREEN)
mtext(s, 0.8, 4.8, 5.5, 2.0, [
    "$ cd cost-analysis-was",
    "$ cp .env.example .env",
    "$ vim .env  # 비밀번호 설정",
    "",
    "$ cd ../cost-analysis-src",
    "$ npm run build",
    "",
    "$ cd ../cost-analysis-was",
    "$ docker-compose up -d",
], 11)

box(s, 6.8, 4.2, 5.8, 2.8, WHITE, RGBColor(200,200,200))
text(s, 7, 4.3, 5.5, 0.4, "접속 URL", 14, True, BLUE)
mtext(s, 7, 4.8, 5.5, 2.0, [
    ("프론트엔드", True, 12, DARK), "  http://localhost",
    ("API 직접", True, 12, DARK), "  http://localhost:8080",
    ("Swagger", True, 12, DARK), "  http://localhost:8080/swagger-ui/index.html",
    ("MinIO Console", True, 12, DARK), "  http://localhost:9001",
], 11)

# Thank You
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), prs.slide_width, Inches(0.06))
b.fill.solid(); b.fill.fore_color.rgb = RED; b.line.fill.background()
text(s, 1.5, 2.2, 10, 1, "Thank You", 48, True, WHITE, PP_ALIGN.CENTER)
text(s, 1.5, 4.0, 10, 0.6, "견적 원가 분석 시스템 — 백엔드 설계서", 20, False, RGBColor(180,180,180), PP_ALIGN.CENTER)

prs.save(r'c:\WORK\cost-analysis-was\docs\백엔드_설계서.pptx')
print("백엔드 설계서 PPTX 생성 완료!")
